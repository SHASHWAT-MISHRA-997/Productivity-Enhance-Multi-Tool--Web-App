import streamlit as st 
import pandas as pd
import PyPDF2
from PyPDF2 import PdfReader, PdfWriter
import pytesseract
from io import BytesIO
from reportlab.pdfgen import canvas
import pdfkit
from rembg import remove
import tempfile
import io
from docx import Document
from pptx import Presentation
import fitz
import zlib
from zipfile import BadZipFile
from PIL import Image, ImageDraw, ImageFont, ImageColor
import os
from fpdf import FPDF
import cv2
import numpy as np
from moviepy.editor import VideoFileClip

def convert_document(input_file, output_format):
    # Dummy implementation for document conversion
    # This is where you'd integrate a library or service for actual conversion
    output_file = io.BytesIO()
    if output_format == 'PDF':
        # For demonstration, just save as PDF
        pdf_writer = PyPDF2.PdfWriter()
        pdf_writer.add_blank_page(width=72, height=72)
        pdf_writer.write(output_file)
    elif output_format == 'DOCX':
        doc = Document()
        doc.add_paragraph("Converted document")
        doc.save(output_file)
    elif output_format == 'PPTX':
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[0])
        slide.shapes.title.text = "Converted presentation"
        presentation.save(output_file)
    output_file.seek(0)
    return output_file

# Video compression function
def compress_video(uploaded_file, compression_level, output_format):
    # Save the uploaded file to a temporary location
    with tempfile.NamedTemporaryFile(delete=False, suffix=".mp4") as temp_file:
        temp_file.write(uploaded_file.read())
        temp_file_path = temp_file.name
    
    # Use the temporary file path with VideoFileClip
    video = VideoFileClip(temp_file_path)
    
    # Create another temporary file to save the compressed video
    compressed_temp_file_path = tempfile.mktemp(suffix=f".{output_format.lower()}")
    
    # Compression logic (write compressed video to the temporary file)
    video.write_videofile(compressed_temp_file_path, codec="libx264", bitrate=f"{compression_level}k")
    video.close()
    
    # Read the compressed video into a BytesIO buffer
    with open(compressed_temp_file_path, "rb") as f:
        compressed_video_bytes = f.read()

    # Clean up temporary files
    os.remove(temp_file_path)
    os.remove(compressed_temp_file_path)

    return io.BytesIO(compressed_video_bytes)


def trim_video(video_file, start_time, end_time):
    try:
        # Save the uploaded video file locally
        temp_file_path = f"temp_{video_file.name}"
        with open(temp_file_path, "wb") as f:
            f.write(video_file.getbuffer())

        # Process the video using moviepy
        with VideoFileClip(temp_file_path) as video:
            trimmed_video = video.subclip(start_time, end_time)
            trimmed_video_path = f"trimmed_{video_file.name}"
            trimmed_video.write_videofile(trimmed_video_path, codec="libx264")
        
        # Remove the temporary file after processing
        os.remove(temp_file_path)
        return trimmed_video_path
    except Exception as e:
        st.error(f"Error processing video {video_file.name}: {e}")
        return None



def merge_pdfs(uploaded_files):
    pdf_writer = PyPDF2.PdfWriter()
    
    for pdf_file in uploaded_files:
        try:
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            for page_num in range(len(pdf_reader.pages)):
                pdf_writer.add_page(pdf_reader.pages[page_num])
        except PyPDF2.errors.PdfReadError as e:
            st.error(f"Error reading {pdf_file.name}: {e}")
    
    output_pdf = io.BytesIO()
    pdf_writer.write(output_pdf)
    output_pdf.seek(0)
    
    # Debugging step: Save to a local file
    with open("debug_merged.pdf", "wb") as f:
        f.write(output_pdf.getvalue())
    
    return output_pdf

def split_pdf(uploaded_file, split_option, custom_ranges=None, split_every=None):
    try:
        reader = PdfReader(uploaded_file)
        writer = PdfWriter()

        if split_option == "Split by Page Ranges":
            ranges = [r.strip() for r in custom_ranges.split(',')]
            for page_range in ranges:
                start, end = map(int, page_range.split('-'))
                for page_num in range(start-1, end):
                    writer.add_page(reader.pages[page_num])

        elif split_option == "Split Every N Pages":
            for i in range(0, len(reader.pages), split_every):
                writer.add_page(reader.pages[i])

        # Save the split PDF to a BytesIO object
        split_pdf_io = BytesIO()
        writer.write(split_pdf_io)
        split_pdf_io.seek(0)  # Reset the buffer position to the beginning
        return split_pdf_io

    except Exception as e:
        st.error(f"An error occurred while splitting the PDF: {e}")
        return None


# Function to compress PDF
def compress_pdf(input_pdf, compression_level):
    # Create a temporary file to save the compressed PDF
    temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=".pdf")
    writer = PdfWriter()

    # Read the uploaded PDF
    pdf_reader = PdfReader(input_pdf)
    for page in pdf_reader.pages:
        writer.add_page(page)

    with open(temp_file.name, "wb") as f:
        writer.write(f)

    # Perform compression (This example uses zlib, but this is a basic approach and might not achieve actual PDF compression)
    with open(temp_file.name, "rb") as f:
        original_data = f.read()
        compressed_data = zlib.compress(original_data, level=compression_level)  # Adjust compression level here

    # Save the compressed data
    compressed_temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=".pdf")
    with open(compressed_temp_file.name, "wb") as f:
        f.write(compressed_data)

    return compressed_temp_file.name


# Function to rotate the PDF
def rotate_pdf(pdf_bytes, rotation_angle, page_selection):
    pdf_document = fitz.open("pdf", pdf_bytes)
    
    # Apply rotation to the selected pages
    if page_selection == "All Pages":
        for page in pdf_document:
            page.set_rotation(rotation_angle)
    else:
        # Select specific pages by index
        for page_number in page_selection:
            if 0 <= page_number < len(pdf_document):
                page = pdf_document[page_number]
                page.set_rotation(rotation_angle)

    # Save the rotated PDF to bytes
    output_pdf_bytes = pdf_document.write()
    return output_pdf_bytes

# Function to convert hex color to RGB tuple with values between 0 and 1
def hex_to_rgb_tuple(hex_color):
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) / 255.0 for i in (0, 2, 4))

# Function to add page numbers to the PDF
def add_page_numbers_to_pdf(pdf_bytes, page_range, position, format_style, font_size, font_color):
    pdf_document = fitz.open("pdf", pdf_bytes)
    total_pages = len(pdf_document)
    
    # Determine which pages to add numbers to
    if page_range == "All Pages":
        pages_to_number = range(total_pages)
    else:
        # Select specific range
        start, end = st.select_slider("Select page range", options=list(range(1, total_pages + 1)), value=(1, total_pages), key="page_range_slider")
        pages_to_number = range(start - 1, end)

    # Mapping positions to coordinates
    position_mapping = {
        "Top-Left": (50, 30),
        "Top-Center": (300, 30),
        "Top-Right": (550, 30),
        "Bottom-Left": (50, 770),
        "Bottom-Center": (300, 770),
        "Bottom-Right": (550, 770),
    }

    # Set the selected position
    pos_x, pos_y = position_mapping[position]

    # Convert font color from hex to RGB tuple
    font_color_rgb = hex_to_rgb_tuple(font_color)

    # Add page numbers to the selected pages
    for page_num in pages_to_number:
        page = pdf_document[page_num]
        text = format_style.format(page_num + 1, total_pages)
        page.insert_text((pos_x, pos_y), text, fontsize=font_size, color=font_color_rgb)

    # Save the modified PDF to bytes
    output_pdf_bytes = pdf_document.write()
    return output_pdf_bytes

# Function to extract text from a PDF using PyMuPDF
def extract_text_from_pdf(pdf_file):
    text = ""
    try:
        pdf_document = fitz.open(stream=pdf_file.read(), filetype="pdf")
        for page_num in range(len(pdf_document)):
            page = pdf_document.load_page(page_num)
            text += page.get_text()
        return text
    except Exception as e:
        st.error(f"An error occurred while extracting text from PDF: {e}")
        return ""

# Function to extract text from an image using OCR
def extract_text_from_image(image):
    text = pytesseract.image_to_string(image)
    return text


# Function to add a text watermark to an image
def add_text_watermark(image, text, position, font_path, font_size, color, opacity, rotation):
    img = image.convert("RGBA")
    watermark_layer = Image.new("RGBA", img.size, (255, 255, 255, 0))
    
    draw = ImageDraw.Draw(watermark_layer)
    font = ImageFont.truetype(font_path, font_size)
    
    # Add text to the watermark layer
    draw.text(position, text, font=font, fill=color + (opacity,))
    
    # Rotate the text layer if needed
    if rotation != 0:
        watermark_layer = watermark_layer.rotate(rotation, expand=1)
    
    # Combine the watermark layer with the original image
    watermarked_image = Image.alpha_composite(img, watermark_layer)
    return watermarked_image

# Function to add an image watermark to an image
def add_image_watermark(image, watermark_image, position, opacity, scale):
    img = image.convert("RGBA")
    watermark = watermark_image.convert("RGBA")
    
    # Scale the watermark image
    width, height = watermark.size
    watermark = watermark.resize((int(width * scale), int(height * scale)), Image.ANTIALIAS)
    
    # Make the watermark transparent
    watermark.putalpha(opacity)
    
    # Add the watermark to the original image
    img.paste(watermark, position, watermark)
    
    return img

def pdf_to_images(pdf_file):
    try:
        pdf_document = fitz.open(stream=pdf_file.read(), filetype="pdf")
        images = []
        
        for page_num in range(len(pdf_document)):
            page = pdf_document.load_page(page_num)
            pix = page.get_pixmap()
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            images.append(img)
        
        return images, pdf_document
    except Exception as e:
        st.error(f"Error processing PDF: {e}")
        return [], None

def detect_and_remove_watermark(image):
    try:
        # Convert image to numpy array for processing
        img_np = np.array(image)
        gray = cv2.cvtColor(img_np, cv2.COLOR_RGB2GRAY)

        # Apply thresholding to highlight potential watermarks
        _, thresh = cv2.threshold(gray, 200, 255, cv2.THRESH_BINARY)

        # Find contours (areas with similar pixel intensities)
        contours, _ = cv2.findContours(thresh, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
        
        # Filter out small contours (likely noise) and fill large contours with white (erase)
        for cnt in contours:
            if cv2.contourArea(cnt) > 100:  # Adjust this threshold based on your watermark size
                x, y, w, h = cv2.boundingRect(cnt)
                cv2.rectangle(img_np, (x, y), (x + w, y + h), (255, 255, 255), -1)  # White fill to erase

        # Convert back to PIL image for displaying
        cleaned_image = image.copy()
         # Example: Draw a rectangle over the watermark area (for demonstration purposes)
        draw = ImageDraw.Draw(cleaned_image)
        draw.rectangle([10, 10, 100, 100], fill="white")  # Modify coordinates and fill color as needed
        
        return cleaned_image
        
    except Exception as e:
        st.error(f"Error detecting or removing watermark: {e}")
        return image  # Return the original image if processing fails

def convert_pil_to_pixmap(pil_image):
    try:
        # Convert PIL image to RGB mode
        pil_image = pil_image.convert("RGB")
        
        # Convert the PIL image to a byte array
        img_byte_arr = io.BytesIO()
        pil_image.save(img_byte_arr, format="PNG")
        img_byte_arr = img_byte_arr.getvalue()
        
        # Convert byte array to fitz Pixmap
        pix = fitz.Pixmap(fitz.open("png", img_byte_arr))
        return pix
    except Exception as e:
        st.error(f"Error converting image for PDF: {e}")
        return None

def apply_watermark_removal(pdf_document, cleaned_images, selected_pages):
    for page_num in selected_pages:
        page = pdf_document.load_page(page_num)
        # Convert cleaned image back to pixmap and replace page content
        pil_image = cleaned_images[page_num]
        pix = convert_pil_to_pixmap(pil_image)
        if pix:
            page.insert_image(page.rect, pixmap=pix)
    
    output_pdf = io.BytesIO()
    pdf_document.save(output_pdf)
    output_pdf.seek(0)
    return output_pdf

def display_image_info(image, label):
    st.text(f"{label} - Size: {image.size}, Mode: {image.mode}")

# Define the images and pages_to_process variables
images = []  # This should be populated with the actual images
pages_to_process = range(len(images))  # Assuming you want to process all images

# Inside your processing loop
cleaned_images = []  # Define the cleaned_images variable
for page_num in pages_to_process:
    st.text(f"Page {page_num + 1}")
    original_image = images[page_num]
    
    if original_image is not None:
        cleaned_image = detect_and_remove_watermark(original_image)
        
        # Display image information for debugging
        display_image_info(original_image, "Original")
        display_image_info(cleaned_image, "Cleaned")
        
        # Display both original and cleaned image side by side for comparison
        col1, col2 = st.columns(2)
        col1.image(original_image, caption=f"Original Page {page_num + 1}")
        col2.image(cleaned_image, caption=f"Cleaned Page {page_num + 1}")

        cleaned_images.append(cleaned_image)
    else:
        st.warning(f"Page {page_num + 1} could not be processed. Image is None.")


# Define passport photo size
PASSPORT_PHOTO_SIZE = (413, 531)  # Example size (width x height) in pixels

def convert_image(image, output_format, quality=None):
    # Convert to RGB if saving as JPEG
    if output_format == "JPEG" and image.mode == "RGBA":
        image = image.convert("RGB")
    
    buffer = io.BytesIO()
    image.save(buffer, format=output_format, quality=quality)
    buffer.seek(0)
    return buffer

def resize_image(image, size, resolution):
    if resolution == "Original":
        return image.resize(size, Image.ANTIALIAS)
    elif resolution == "Full HD":
        return image.resize((1920, 1080), Image.ANTIALIAS)
    elif resolution == "HD":
        return image.resize((1280, 720), Image.ANTIALIAS)
    elif resolution == "SD":
        return image.resize((720, 480), Image.ANTIALIAS)
    else:
        raise ValueError("Invalid resolution")

def compress_image(uploaded_image, compression_level, quality, resolution):
    # Convert uploaded image to RGB mode
    img = Image.open(uploaded_image).convert("RGB")
    
    # Define quality levels
    quality_settings = {
        "Low": 30,
        "Medium": 60,
        "High": 90
    }
    quality = quality_settings.get(quality, 90)
    
    # Resize image to the chosen resolution
    if resolution == "Full HD":
        img = img.resize((1920, 1080), Image.Resampling.LANCZOS)
    elif resolution == "HD":
        img = img.resize((1280, 720), Image.Resampling.LANCZOS)
    elif resolution == "SD":
        img = img.resize((640, 480), Image.Resampling.LANCZOS)
    
    # Adjust compression level
    if compression_level == "High":
        quality = 50
    elif compression_level == "Medium":
        quality = 70
    elif compression_level == "Low":
        quality = 90

    # Compress image
    buffer = io.BytesIO()
    img.save(buffer, format="JPEG", quality=quality)
    buffer.seek(0)
    compressed_image = Image.open(buffer)
    
    # Calculate the size of the compressed image
    compressed_image_size_kb = len(buffer.getvalue()) / 1024
    
    return buffer, quality, compressed_image_size_kb

# Function to remove the background and replace it with a chosen color and transparency
def remove_background(image, background_color, transparency):
    image = Image.open(image)
    image = remove(image)  # Remove background

    # Convert hex color to RGBA with transparency
    rgba_color = tuple(int(background_color[i:i + 2], 16) for i in (1, 3, 5)) + (int(transparency * 255),)
    background = Image.new("RGBA", image.size, rgba_color)  # Create background with chosen color and transparency

    image = image.convert("RGBA")  # Ensure image has an alpha channel
    final_image = Image.alpha_composite(background, image)  # Composite the image onto the background

    return final_image  # Return the image in RGBA format to preserve transparency


def passport_size_photo(image):
    try:
        img = Image.open(image).convert("RGB")  # Convert to RGB mode
        passport_photo = img.resize(PASSPORT_PHOTO_SIZE, Image.LANCZOS)
        buffer = BytesIO()
        passport_photo.save(buffer, format="JPEG")
        buffer.seek(0)
        return buffer
    except Exception as e:
        st.error(f"Error in passport photo creation: {e}")
        return None

def print_passport_photo(image, page_size):
    try:
        img = Image.open(image).convert("RGB")  # Convert to RGB mode
        passport_photo = img.resize(PASSPORT_PHOTO_SIZE, Image.LANCZOS)
        page_width, page_height = page_size
        page_image = Image.new('RGB', (page_width, page_height), color='white')
        photo_width, photo_height = passport_photo.size
        position = ((page_width - photo_width) // 2, (page_height - photo_height) // 2)
        page_image.paste(passport_photo, position)
        buffer = BytesIO()
        page_image.save(buffer, format="JPEG")
        buffer.seek(0)
        return buffer
    except Exception as e:
        st.error(f"Error in creating passport photo page: {e}")
        return None
    
def add_text_watermark(pdf_file, text, font_size, color, opacity, rotation, position):
    doc = fitz.open(stream=pdf_file.read(), filetype="pdf")
    
    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        rect = page.rect
        pil_img = Image.new("RGBA", (int(rect.width), int(rect.height)), (255, 255, 255, 0))
        draw = ImageDraw.Draw(pil_img)
        
        try:
            font = ImageFont.truetype("arial.ttf", font_size)
        except IOError:
            font = ImageFont.load_default()

        # Use textbbox to calculate text size
        text_bbox = draw.textbbox((0, 0), text, font=font)
        text_width = text_bbox[2] - text_bbox[0]
        text_height = text_bbox[3] - text_bbox[1]

        text_x, text_y = {
            "Top-Left": (10, 10),
            "Top-Right": (rect.width - text_width - 10, 10),
            "Center": ((rect.width - text_width) // 2, (rect.height - text_height) // 2),
            "Bottom-Left": (10, rect.height - text_height - 10),
            "Bottom-Right": (rect.width - text_width - 10, rect.height - text_height - 10)
        }[position]

        # Apply rotation
        watermark_img = pil_img.rotate(rotation, expand=True)
        draw = ImageDraw.Draw(watermark_img)
        draw.text((text_x, text_y), text, font=font, fill=(ImageColor.getrgb(color) + (opacity,)))
        
        img_bytes = io.BytesIO()
        watermark_img.save(img_bytes, format='PNG')
        img_bytes.seek(0)
        page.insert_image(rect, stream=img_bytes.read(), overlay=True)
        
    output_pdf = io.BytesIO()
    doc.save(output_pdf)
    output_pdf.seek(0)
    return output_pdf
def add_image_watermark(pdf_file, watermark_file, opacity, scale, position):
    doc = fitz.open(stream=pdf_file.read(), filetype="pdf")
    watermark_img = Image.open(watermark_file).convert("RGBA")
    watermark_img = watermark_img.resize((int(watermark_img.width * scale), int(watermark_img.height * scale)))
    watermark_img.putalpha(opacity)
    
    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        rect = page.rect
        watermark_pil = Image.new("RGBA", (int(rect.width), int(rect.height)))
        watermark_pil.paste(watermark_img, {
            "Top-Left": (10, 10),
            "Top-Right": (rect.width - watermark_img.width - 10, 10),
            "Center": ((rect.width - watermark_img.width) // 2, (rect.height - watermark_img.height) // 2),
            "Bottom-Left": (10, rect.height - watermark_img.height - 10),
            "Bottom-Right": (rect.width - watermark_img.width - 10, rect.height - watermark_img.height - 10)
        }[position], watermark_img)
        
        img_bytes = io.BytesIO()
        watermark_pil.save(img_bytes, format='PNG')
        img_bytes.seek(0)
        page.insert_image(rect, stream=img_bytes.read(), overlay=True)
        
    output_pdf = io.BytesIO()
    doc.save(output_pdf)
    output_pdf.seek(0)
    return output_pdf


def edit_pdf(input_pdf):
    try:
        reader = PdfReader(input_pdf)
        writer = PdfWriter()

        # Implement your PDF editing logic here
        for page in reader.pages:
            writer.add_page(page)

        output = BytesIO()
        writer.write(output)
        output.seek(0)
        return output.getvalue()  # Return the binary data
    except Exception as e:
        print(f"Error in edit_pdf: {e}")
        return None
    

def add_text_to_pdf(pdf_path, text, x, y):
    reader = PdfReader(pdf_path)
    writer = PdfWriter()

    # Create a temporary PDF with the text to be added
    for page in reader.pages:
        packet = io.BytesIO()
        can = FPDF()
        can.add_page()
        can.set_xy(x, y)
        can.set_font("Arial", size=12)
        can.multi_cell(0, 10, text)

        # Save the FPDF content into the BytesIO stream
        can.output(packet)
        packet.seek(0)

        # Merge the text PDF with the existing page
        text_pdf = PdfReader(packet)
        page.merge_page(text_pdf.pages[0])
        writer.add_page(page)

    # Write the final output to a BytesIO object
    output_stream = io.BytesIO()
    writer.write(output_stream)
    output_stream.seek(0)  # Reset the stream position
    return output_stream

def add_image_to_pdf(pdf_path, image_path, x, y, width, height):
    reader = PdfReader(pdf_path)
    writer = PdfWriter()

    for page in reader.pages:
        packet = io.BytesIO()

        # Load the image and resize it
        img = Image.open(image_path)
        img = img.resize((width, height))

        # Create a temporary PDF with the image
        img_pdf = FPDF()
        img_pdf.add_page()
        img.save(packet, format="PNG")
        packet.seek(0)
        img_pdf.image(packet, x=x, y=y, w=width, h=height)

        # Save the FPDF content into the BytesIO stream
        img_output = io.BytesIO()
        img_pdf.output(img_output)
        img_output.seek(0)

        # Merge the image PDF with the existing page
        image_pdf = PdfReader(img_output)
        page.merge_page(image_pdf.pages[0])
        writer.add_page(page)

    # Write the final output to a BytesIO object
    output_stream = io.BytesIO()
    writer.write(output_stream)
    output_stream.seek(0)  # Reset the stream position
    return output_stream


# Set page configuration as the very first command in your script
st.set_page_config(page_title="Advanced PDF Editor", layout="wide")

def crop_pdf(pdf_path, crop_box):
    reader = PdfReader(pdf_path)
    writer = PdfWriter()

    for page in reader.pages:
        page.cropbox.lower_left = (crop_box[0], crop_box[1])
        page.cropbox.upper_right = (crop_box[2], crop_box[3])
        writer.add_page(page)

    output_stream = io.BytesIO()
    writer.write(output_stream)
    output_stream.seek(0)
    return output_stream

def draw_on_pdf(pdf_path, shape_type, color, dimensions):
    reader = PdfReader(pdf_path)
    writer = PdfWriter()

    for page in reader.pages:
        packet = io.BytesIO()
        can = canvas.Canvas(packet)
        can.setStrokeColorRGB(*color)
        can.setFillColorRGB(*color)

        if shape_type == "Rectangle":
            can.rect(*dimensions)
        elif shape_type == "Circle":
            can.circle(dimensions[0], dimensions[1], dimensions[2])
        elif shape_type == "Line":
            can.line(*dimensions)

        can.save()
        packet.seek(0)

        overlay_pdf = PdfReader(packet)
        page.merge_page(overlay_pdf.pages[0])
        writer.add_page(page)

    output_stream = io.BytesIO()
    writer.write(output_stream)
    output_stream.seek(0)
    return output_stream

def highlight_text_in_pdf(pdf_path, color, dimensions):
    reader = PdfReader(pdf_path)
    writer = PdfWriter()

    for page in reader.pages:
        packet = io.BytesIO()
        can = canvas.Canvas(packet)
        can.setFillColorRGB(*color)
        can.rect(*dimensions, stroke=0, fill=1)
        can.save()
        packet.seek(0)

        overlay_pdf = PdfReader(packet)
        page.merge_page(overlay_pdf.pages[0])
        writer.add_page(page)

    output_stream = io.BytesIO()
    writer.write(output_stream)
    output_stream.seek(0)
    return output_stream

def zoom_pdf(pdf_path, zoom_level):
    reader = PdfReader(pdf_path)
    writer = PdfWriter()

    for page in reader.pages:
        page.scale_by(zoom_level)
        writer.add_page(page)

    output_stream = io.BytesIO()
    writer.write(output_stream)
    output_stream.seek(0)
    return output_stream

def add_link_to_pdf(pdf_path, link_url, rect):
    reader = PdfReader(pdf_path)
    writer = PdfWriter()

    for page in reader.pages:
        page.add_uri(rect, link_url)
        writer.add_page(page)

    output_stream = io.BytesIO()
    writer.write(output_stream)
    output_stream.seek(0)
    return output_stream

def html_to_pdf(html_content):
    try:
        # Convert HTML to PDF using pdfkit
        pdf_output = pdfkit.from_string(html_content, False)  # Returns a binary string
        return pdf_output
    except Exception as e:
        st.error(f"Failed to convert HTML to PDF: {e}")
        return None

def unlock_pdf(uploaded_file):
    # Placeholder function to unlock a protected PDF
    pass



def protect_pdf(input_pdf, password):
    try:
        reader = PdfReader(input_pdf)
        writer = PdfWriter()

        for page in reader.pages:
            writer.add_page(page)

        writer.encrypt(password)
        
        # Save the protected PDF to a BytesIO object
    
        output = BytesIO()
        writer.write(output)
        output.seek(0)
        return output.getvalue()  # Return the binary data
    except Exception as e:
        print(f"Error in protect_pdf: {e}")
        return None
def reorganize_pdf(uploaded_pdf):
    try:
        # Example logic to reorganize PDF (e.g., rearrange pages)
        reader = PdfReader(uploaded_pdf)
        writer = PdfWriter()

        # Reorganize pages (e.g., reverse order)
        for page in reversed(reader.pages):
            writer.add_page(page)

        # Save reorganized PDF to a buffer
        pdf_buffer = io.BytesIO()
        writer.write(pdf_buffer)
        pdf_buffer.seek(0)
        return pdf_buffer.read()  # Return binary content
    except Exception as e:
        st.error(f"Error while reorganizing PDF: {e}")
        return None
def repair_pdf(uploaded_pdf):
    try:
        # Example repair logic: Just reading and writing back (basic check)
        reader = PdfReader(uploaded_pdf)
        writer = PdfWriter()

        for page in reader.pages:
            writer.add_page(page)

        # Save the repaired PDF to a buffer
        pdf_buffer = io.BytesIO()
        writer.write(pdf_buffer)
        pdf_buffer.seek(0)
        return pdf_buffer.read()  # Return binary content
    except Exception as e:
        st.error(f"Error while repairing PDF: {e}")
        return None


# Helper Functions# Function to convert PDF to JPG
def convert_pdf_to_jpg(pdf_file):
    pdf = fitz.open(stream=pdf_file.read(), filetype="pdf")
    images = []
    for page_num in range(len(pdf)):
        page = pdf.load_page(page_num)
        pix = page.get_pixmap()
        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
        images.append(img)
    return images

# Function to convert JPG to PDF
def convert_jpg_to_pdf(image_files):
    pdf_bytes = io.BytesIO()
    images = [Image.open(image_file).convert("RGB") for image_file in image_files]
    if images:
        images[0].save(pdf_bytes, format="PDF", save_all=True, append_images=images[1:], resolution=100.0, quality=95)
        pdf_bytes.seek(0)
    return pdf_bytes


def convert_excel_to_pdf(excel_file):
    df = pd.read_excel(excel_file)
    buffer = io.BytesIO()
    with pd.ExcelWriter(buffer, engine="xlsxwriter") as writer:
        df.to_excel(writer, sheet_name="Sheet1", index=False)
    buffer.seek(0)
    return buffer


def convert_pdf_to_excel(pdf_file):
    # Placeholder: Implement PDF to Excel logic.
    # For now, just return some dummy content.
    output = io.BytesIO()
    df = pd.DataFrame({"Sample Data": [1, 2, 3, 4, 5]})
    with pd.ExcelWriter(output, engine="xlsxwriter") as writer:
        df.to_excel(writer, sheet_name="Sheet1", index=False)
    output.seek(0)
    return output


def convert_powerpoint_to_pdf(uploaded_file):
    try:
        if uploaded_file is not None and uploaded_file.name.endswith('.pptx'):
            # Simulate the conversion process
            pdf_data = b"Fake PDF binary data"
            
            # Normally, you would return the actual PDF binary data here
            return pdf_data
        else:
            st.error("Uploaded file is not a valid .pptx PowerPoint file.")
            return None
    except Exception as e:
        st.error(f"An error occurred during conversion: {e}")
        return None

def convert_pdf_to_powerpoint(pdf_file):
    # Placeholder: Implement PDF to PowerPoint logic.
    return "PDF to PowerPoint conversion is in progress."
def convert_word_to_pdf(docx_file):
    # Example conversion logic
    document = Document(docx_file)
    pdf_buffer = BytesIO()

    # Here we would convert the DOCX to PDF using an appropriate library/method
    # For example, you can use 'pdfkit' with 'wkhtmltopdf', or any other method
    # Assuming you write the PDF bytes to 'pdf_buffer'

    return pdf_buffer.getvalue()
def convert_pdf_to_word(pdf_file):
    # Simulate the conversion process (replace with actual conversion logic)
    if pdf_file is not None:
        # Normally, this function would return the actual Word document's binary data
        word_data = b"Fake Word binary data"
        return word_data
    return None

def display_pdf_with_pdfjs(pdf_path):
    pdf_url = f"data:application/pdf;base64,{pdf_path}"
    html = f"""
    <iframe src="https://mozilla.github.io/pdf.js/web/viewer.html?file={pdf_url}" width="100%" height="800px">
    </iframe>
    """
    st.components.v1.html(html, height=800)

def display_pdf_as_images(pdf_path):
    try:
        # Open the PDF file
        pdf_document = fitz.open(pdf_path)
        for page_num in range(len(pdf_document)):
            # Load the page
            page = pdf_document.load_page(page_num)
            
            # Render the page as an image
            pix = page.get_pixmap()
            
            # Convert the image to PIL format
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            
            # Display the image in Streamlit
            buffer = io.BytesIO()
            img.save(buffer, format="PNG")
            st.image(buffer, caption=f"Page {page_num + 1}", use_column_width=True)
    except Exception as e:
        st.error(f"An error occurred: {e}")

def unlock_pdf_with_brute_force(uploaded_pdf):
    # Common password dictionary (you can expand this list)
    common_passwords = [
        "1234", "password", "12345", "123456", "password123", "admin", "secret", "letmein"
        # Add more passwords or load from a file if needed
    ]

    # Create a PDF reader object
    reader = PyPDF2.PdfReader(uploaded_pdf)
    
    if not reader.is_encrypted:
        st.success("This PDF is not password-protected.")
        return None

    # Attempt to decrypt using each password
    for password in common_passwords:
        try:
            if reader.decrypt(password) == 1:  # Decrypt returns 1 on success
                st.success(f"Password found: '{password}'")
                # Create a PDF writer object for the unlocked PDF
                writer = PyPDF2.PdfWriter()

                # Add all pages to the writer object
                for page in range(len(reader.pages)):
                    writer.add_page(reader.pages[page])

                # Save the unlocked PDF to a BytesIO stream
                unlocked_pdf_stream = io.BytesIO()
                writer.write(unlocked_pdf_stream)
                unlocked_pdf_stream.seek(0)  # Reset stream position for downloading

                return unlocked_pdf_stream
        except:
            continue

    st.error("Password not found in the dictionary. Try another method.")
    return None


def main():
    

    # Sidebar
    st.sidebar.title("Navigation")
    options = ["Home", "About Developer", "Document Features", "Image Features", "Video Features"]
    choice = st.sidebar.radio("Go to", options)



    # Home Page
   
    if choice == "Home":
        st.title("Welcome to the Productivity-Multi-Tool App")
        
        st.header("About This App")
        st.write("""
        **App Name:** Productivity-Enhance-Multi-Tool App  
        **Description:** This app offers a comprehensive suite of tools for handling PDFs, images, videos, and various document formats. With features like compression, conversion, editing, and more, it aims to enhance your productivity.
        """)
        
        st.subheader("Procedure to Use the App:")
        st.write("""
        - **→ Document Features:** 
            - **Merge PDFs:** Combine multiple PDF files into a single document.
            - **Split PDFs:** Divide large PDFs into smaller sections based on page ranges or custom splits.
            - **Compress PDFs:** Reduce the file size of your PDFs while maintaining quality.
            - **Rotate PDFs:** Rotate pages within a PDF to the desired orientation.
            - **Add Page Numbers:** Insert page numbers in your PDFs at customizable positions.
            - **Extract Text:** Pull text content from PDFs using advanced OCR technology.
            - **Add Watermarks:** Secure your documents by adding text or image watermarks.
            - **Remove Watermarks:** Automatically detect and remove watermarks from PDF documents.
            - **Edit PDFs:** Crop, draw, highlight text, zoom, or add links to your PDF files.
            - **Convert PDFs to Images:** Convert PDF pages to high-quality images (JPG/PNG).
            - **Convert Images to PDFs:** Combine multiple images into a single PDF file.
            - **Convert PDFs to Excel:** Transform PDF tables into editable Excel sheets.
            - **Convert Excel to PDF:** Convert Excel files to PDF format for easy sharing.
            - **Convert PowerPoint to PDF:** Convert PowerPoint presentations into PDF documents.
            - **Convert PDFs to PowerPoint:** Transform PDFs into PowerPoint slides.
            - **Convert Word to PDF:** Convert Word documents to PDF format.
            - **Convert PDFs to Word:** Extract text and images from PDFs into editable Word documents.
            - **Unlock PDFs:** Remove passwords from protected PDFs.
            - **Protect PDFs:** Add password protection to your PDFs.
            - **Repair PDFs:** Fix corrupted PDF files and restore them to a usable state.
            - **Reorganize PDFs:** Rearrange, reorder, or reverse the pages of your PDF documents.
        - **→ Image Features:** 
            - Convert images between formats (JPEG, PNG, etc.).
            - Compress images to reduce file size.
            - Remove backgrounds from images for a cleaner look.
            - Resize images or create passport-sized photos.
        - **→ Video Compressor Features:** 
            - Compress videos to reduce their size.
            - Trim specific parts of a video.
        """)
        
        st.header("Developed by Shashwat Mishra")
        st.write("""
        Hi Everyone! I'm Shashwat Mishra, a Robotics and Automation graduate. This app is part of my continuous effort to develop efficient and user-friendly tools. I hope this app helps you streamline your tasks!
        """)

        

    # About Developer Page
    elif choice == "About Developer":
        st.title("About the Developer")
        st.write("""
            Hi Everyone, I am Shashwat Mishra. I am a 2024 BTech Passout Student in Robotics and Automation. I have a keen interest in developing efficient and effective solutions. I have worked on various projects.
            
            My primary focus is on creating user-friendly and efficient tools that can help users improve their productivity and efficiency. I am always looking for new ideas and opportunities to improve my skills and knowledge.
            
            Here is my Another Project That I have Created : 
                https://whatsapp-chat-advanced-ai-analyzer-web-app-shashwat-mishra-997.streamlit.app/
            
                Here is my Linkedin Profile For Rest of Project that I have made till Now:
                 
            [LinkedIn Profile](https://www.linkedin.com/in/sm980)  """)

           

    if choice == "Document Features":
        st.title("Document Features")

        st.title("Welcome to the Advanced-Document Editor :")

        st.header("Editing Tools")
        st.write("Select an editing tool below:")
        option = st.selectbox("Choose an Action", ( "PDF to JPG","JPG to PDF","Excel to PDF","PDF to Excel","PowerPoint to PDF","PDF to PowerPoint","Word to PDF",
        "PDF to Word","Merge PDFs", "Split PDF", "Compress PDF", "Rotate PDF", "Add Page Numbers", "Extract Text", "Add Watermark", "Remove Watermark", "Edit PDF", "HTML to PDF", "Unlock PDF","Protect PDF",
        "Organize PDF", "Repair PDF" ))



        if option == "Merge PDFs":        
            st.subheader("Merge PDFs")
            uploaded_files = st.file_uploader("Upload PDF files to merge", type="pdf", accept_multiple_files=True)
            if st.button("Merge PDFs"):
                if uploaded_files:
                    merged_pdf_io = merge_pdfs(uploaded_files)
                    st.download_button(
                        label="Download Merged PDF",
                        data=merged_pdf_io,
                        file_name="merged.pdf",
                        mime="application/pdf"
                    )
                else:
                    st.warning("Please upload PDF files to merge.")


        elif option == "Split PDF":
            st.subheader("Split PDF")
            uploaded_file = st.file_uploader("Upload PDF to split", type=["pdf"], key="split_pdf")
            split_option = st.selectbox("Split Option", ["Split by Page Ranges", "Split Every N Pages"])

            custom_ranges = None
            split_every = None

            if split_option == "Split by Page Ranges":
                custom_ranges = st.text_input("Enter page ranges (e.g., 1-3, 4-6)")
            elif split_option == "Split Every N Pages":
                split_every = st.number_input("Enter the number of pages per split", min_value=1, step=1)

            if st.button("Split PDF"):
                if uploaded_file:
                    split_pdf_io = split_pdf(uploaded_file, split_option, custom_ranges, split_every)

                    if split_pdf_io:
                        st.download_button(
                            "Download Split PDF",
                            split_pdf_io,
                            "split_pdf.pdf",
                            mime="application/pdf"
                        )
                    else:
                        st.error("Failed to generate the split PDF. Please check your input and try again.")
                else:
                    st.warning("Please upload a PDF to split.")


        elif option == "Compress PDF":
            st.subheader("Compress PDF")
            # File uploader
            uploaded_file = st.file_uploader("Choose a PDF file to compress", type="pdf")

            # Compression options
            compression_level = st.selectbox(
                "Select Compression Level",
                ["Low", "Medium", "High"],
                help="Choose the level of compression. Higher levels reduce file size more but may reduce quality."
            )

            # Handle file upload and compression
            if uploaded_file:
                st.write("File uploaded successfully!")

                # Compression level mapping
                level_mapping = {
                    "Low": 1,
                    "Medium": 5,
                    "High": 9
                }
                compression_ratio = level_mapping[compression_level]

                if st.button("Compress PDF"):
                    with st.spinner("Compressing PDF..."):
                        compressed_pdf_path = compress_pdf(uploaded_file, compression_ratio)
                        
                        with open(compressed_pdf_path, "rb") as f:
                            st.download_button(
                                label="Download Compressed PDF",
                                data=f,
                                file_name="compressed.pdf",
                                mime="application/pdf"
                            )
                        
                        st.success("Compression completed!")

            # Information and help section
            st.info("This app allows you to compress your PDF files. Choose the file you want to compress, select the desired compression level, and download the compressed PDF.")

        elif option == "Rotate PDF":
            st.subheader("Rotate PDF")
            # Upload PDF
            uploaded_file = st.file_uploader("Choose a PDF file", type="pdf")
            
            if uploaded_file:
                pdf_bytes = uploaded_file.read()
                
                # Display rotation options
                st.subheader("Rotation Settings")
                rotation_angle = st.selectbox("Select rotation angle", [0, 90, 180, 270])
                
                # Display page selection options
                page_count = len(fitz.open("pdf", pdf_bytes))
                page_selection = st.radio(
                    "Select pages to rotate",
                    ["All Pages", "Specific Pages"]
                )
                
                specific_pages = []
                if page_selection == "Specific Pages":
                    specific_pages = st.multiselect(
                        f"Select pages (1-{page_count})",
                        list(range(1, page_count + 1)),
                        format_func=lambda x: f"Page {x}"
                    )
                    specific_pages = [p - 1 for p in specific_pages]  # Adjust to zero-indexing

                # Rotate PDF
                if st.button("Rotate PDF"):
                    rotated_pdf = rotate_pdf(pdf_bytes, rotation_angle, specific_pages if page_selection == "Specific Pages" else "All Pages")
                    st.success("PDF rotated successfully!")

                    # Download rotated PDF
                    st.download_button(
                        label="Download Rotated PDF",
                        data=rotated_pdf,
                        file_name="rotated_pdf.pdf",
                        mime="application/pdf"
                    )
        elif option == "Add Page Numbers":
            st.subheader("Add Page Numbers")
            # Upload PDF
            uploaded_file = st.file_uploader("Choose a PDF file", type="pdf", key="pdf_file_uploader")
            
            if uploaded_file:
                pdf_bytes = uploaded_file.read()
                
                # Page range options
                st.subheader("Page Numbering Settings")
                page_range = st.radio("Select pages to number", ["All Pages", "Custom Range"], key="page_range_radio")

                # Position options
                position = st.selectbox(
                    "Select position for page numbers",
                    ["Top-Left", "Top-Center", "Top-Right", "Bottom-Left", "Bottom-Center", "Bottom-Right"],
                    key="position_selectbox"
                )

                # Numbering format
                format_style = st.selectbox(
                    "Select numbering format",
                    ["Page {}/{}", "Page {}", "{}", "{}/{}"],
                    key="format_selectbox"
                )

                # Font options
                font_size = st.slider("Select font size", 8, 24, 12, key="font_size_slider")
                font_color = st.color_picker("Pick a font color", "#000000", key="font_color_picker")

                # Add page numbers
                if st.button("Add Page Numbers", key="add_page_numbers_button"):
                    result_pdf = add_page_numbers_to_pdf(pdf_bytes, page_range, position, format_style, font_size, font_color)
                    st.success("Page numbers added successfully!")

                    # Download the PDF
                    st.download_button(
                        label="Download PDF with Page Numbers",
                        data=result_pdf,
                        file_name="pdf_with_page_numbers.pdf",
                        mime="application/pdf",
                        key="download_button"
                    )

        elif option == "Extract Text":
            st.subheader("Extract Text from PDFs")
            uploaded_file = st.file_uploader("Upload a PDF file to extract text", type="pdf")
            if st.button("Extract Text"):
                if uploaded_file:
                    text = extract_text_from_pdf(uploaded_file)
                    st.text_area("Extracted Text", text)
                else:
                    st.warning("Please upload a PDF file.")
        
        elif option == "Add Watermark":
            st.title("PDF Watermarking Tool")

            uploaded_pdf = st.file_uploader("Upload PDF", type="pdf")
            watermark_type = st.selectbox("Select Watermark Type", ["Text", "Image"])
            
            if watermark_type == "Text":
                text = st.text_input("Enter Watermark Text")
                font_size = st.slider("Font Size", 10, 200, 50)
                color = st.color_picker("Text Color", "#FFFFFF")
                opacity = st.slider("Opacity", 0, 255, 128)
                rotation = st.slider("Rotation", 0, 360, 0)
                position = st.selectbox("Select Watermark Position", ["Top-Left", "Top-Right", "Center", "Bottom-Left", "Bottom-Right"])
                
                if st.button("Apply Text Watermark"):
                    if uploaded_pdf:
                        result_pdf = add_text_watermark(uploaded_pdf, text, font_size, color, opacity, rotation, position)
                        st.download_button("Download Watermarked PDF", result_pdf, file_name="watermarked.pdf", mime="application/pdf")
                    else:
                        st.error("Please upload a PDF file.")
            
            elif watermark_type == "Image":
                watermark_image = st.file_uploader("Upload Watermark Image", type=["png", "jpg", "jpeg"])
                opacity = st.slider("Opacity", 0, 255, 128)
                scale = st.slider("Scale", 0.1, 1.0, 0.5)
                position = st.selectbox("Select Watermark Position", ["Top-Left", "Top-Right", "Center", "Bottom-Left", "Bottom-Right"])
                
                if st.button("Apply Image Watermark"):
                    if uploaded_pdf and watermark_image:
                        result_pdf = add_image_watermark(uploaded_pdf, watermark_image, opacity, scale, position)
                        st.download_button("Download Watermarked PDF", result_pdf, file_name="watermarked.pdf", mime="application/pdf")
                    else:
                        st.error("Please upload a PDF file and an image file.")



        elif option == "Remove Watermark":
            st.subheader("Remove Watermark")
            uploaded_file = st.file_uploader("Upload a PDF file to remove watermarks", type="pdf", key="unique_file_uploader")

            if uploaded_file:
                images, pdf_document = pdf_to_images(uploaded_file)
                
                if images and pdf_document:
                    st.subheader("Select Pages to Remove Watermark")
                    total_pages = len(images)
                    page_options = ["All Pages"] + [f"Page {i + 1}" for i in range(total_pages)]
                    selected_pages = st.multiselect(
                        "Select Pages",
                        options=page_options,
                        default=["All Pages"],
                        key="page_selector"
                    )

                    # Determine the pages to be processed
                    if "All Pages" in selected_pages:
                        pages_to_process = list(range(total_pages))
                    else:
                        pages_to_process = [int(page.split(" ")[1]) - 1 for page in selected_pages]

                    st.subheader("Detected Watermarks and Erasing")
                    cleaned_images = []

                    # Process each selected page
                    for page_num in pages_to_process:
                        st.text(f"Processing Page {page_num + 1}")
                        original_image = images[page_num]
                        
                        if original_image is not None:
                            cleaned_image = detect_and_remove_watermark(original_image)
                            
                            # Debug: Show image size and mode
                            st.text(f"Original Page {page_num + 1} - Size: {original_image.size}, Mode: {original_image.mode}")
                            st.text(f"Cleaned Page {page_num + 1} - Size: {cleaned_image.size}, Mode: {cleaned_image.mode}")
                            
                            # Display both original and cleaned image side by side for comparison
                            col1, col2 = st.columns(2)
                            col1.image(original_image, caption=f"Original Page {page_num + 1}")
                            col2.image(cleaned_image, caption=f"Cleaned Page {page_num + 1}")

                            cleaned_images.append(cleaned_image)
                        else:
                            st.warning(f"Page {page_num + 1} could not be processed. Image is None.")

                    if st.button("Remove Watermark and Download PDF", key="remove_watermark_button"):
                        if cleaned_images:
                            cleaned_pdf = apply_watermark_removal(pdf_document, cleaned_images, pages_to_process)
                            st.download_button(
                                label="Download Cleaned PDF",
                                data=cleaned_pdf,
                                file_name="cleaned.pdf",
                                mime="application/pdf",
                                key="download_button"
                            )
                        else:
                            st.warning("No images processed. Please check the uploaded file and try again.")
                else:
                    st.warning("Could not process the PDF. Please try a different file.")
            else:
                st.warning("Please upload a PDF file.")

        # Tab for Editing PDF
        elif option == "Edit PDF":
            st.header("Edit PDF")
            action = st.selectbox("Choose an Action", (
                "Crop PDF", "Draw on PDF", "Highlight Text", "Zoom In/Out", "Add Link"
            ))

            uploaded_pdf = st.file_uploader("Upload a PDF file", type="pdf", key="pdf_upload")
            if uploaded_pdf is not None:
                # Save uploaded PDF
                with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp_file:
                    tmp_file.write(uploaded_pdf.getbuffer())
                    pdf_path = tmp_file.name

                st.success("PDF uploaded successfully!")

                # Display PDF preview
                display_pdf_with_pdfjs(pdf_path)

                if action == "Crop PDF":
                    st.subheader("Crop PDF")
                    left = st.number_input("Left", min_value=0)
                    bottom = st.number_input("Bottom", min_value=0)
                    right = st.number_input("Right", min_value=0)
                    top = st.number_input("Top", min_value=0)
                    if st.button("Apply Crop"):
                        crop_box = (left, bottom, right, top)
                        edited_pdf = crop_pdf("temp.pdf", crop_box)
                        st.download_button("Download Cropped PDF", edited_pdf, file_name="cropped.pdf")

                elif action == "Draw on PDF":
                    st.subheader("Draw on PDF")
                    shape_type = st.selectbox("Shape Type", ("Rectangle", "Circle", "Line"))
                    color = st.color_picker("Pick a Color")
                    color_rgb = tuple(int(color[i:i+2], 16)/255 for i in (1, 3, 5))
                    if shape_type == "Rectangle" or shape_type == "Line":
                        dimensions = (
                            st.number_input("X1", min_value=0),
                            st.number_input("Y1", min_value=0),
                            st.number_input("X2", min_value=0),
                            st.number_input("Y2", min_value=0),
                        )
                    elif shape_type == "Circle":
                        dimensions = (
                            st.number_input("X", min_value=0),
                            st.number_input("Y", min_value=0),
                            st.number_input("Radius", min_value=0),
                        )
                    if st.button("Apply Drawing"):
                        edited_pdf = draw_on_pdf("temp.pdf", shape_type, color_rgb, dimensions)
                        st.download_button("Download PDF with Drawing", edited_pdf, file_name="drawing.pdf")

                elif action == "Highlight Text":
                    st.subheader("Highlight Text")
                    color = st.color_picker("Pick Highlight Color")
                    color_rgb = tuple(int(color[i:i+2], 16)/255 for i in (1, 3, 5))
                    dimensions = (
                        st.number_input("X1", min_value=0),
                        st.number_input("Y1", min_value=0),
                        st.number_input("Width", min_value=0),
                        st.number_input("Height", min_value=0),
                    )
                    if st.button("Apply Highlight"):
                        edited_pdf = highlight_text_in_pdf("temp.pdf", color_rgb, dimensions)
                        st.download_button("Download Highlighted PDF", edited_pdf, file_name="highlighted.pdf")

                elif action == "Zoom In/Out":
                    st.subheader("Zoom In/Out")
                    zoom_level = st.slider("Zoom Level", 0.5, 3.0, 1.0, step=0.1)
                    if st.button("Apply Zoom"):
                        edited_pdf = zoom_pdf("temp.pdf", zoom_level)
                        st.download_button("Download Zoomed PDF", edited_pdf, file_name="zoomed.pdf")

                elif action == "Add Link":
                    st.subheader("Add Link to PDF")
                    link_url = st.text_input("Enter the URL")
                    x1 = st.number_input("X1", min_value=0)
                    y1 = st.number_input("Y1", min_value=0)
                    x2 = st.number_input("X2", min_value=0)
                    y2 = st.number_input("Y2", min_value=0)
                    rect = (x1, y1, x2, y2)
                    if st.button("Add Link"):
                        edited_pdf = add_link_to_pdf("temp.pdf", link_url, rect)
                        st.download_button("Download PDF with Link", edited_pdf, file_name="linked.pdf")
        
        # Tab for HTML to PDF
        elif option == "HTML to PDF":
            st.title("HTML to PDF Converter")

            # HTML input
            html_content = st.text_area("Enter HTML content here", height=300, value="<h1>Hello, PDF!</h1>")

            if st.button("Convert to PDF"):
                pdf_content = html_to_pdf(html_content)
                
                if pdf_content:
                    # Display download button for the generated PDF
                    st.download_button(
                        label="Download PDF",
                        data=pdf_content,
                        file_name="converted.pdf",
                        mime="application/pdf"
                    )
                else:
                    st.error("PDF generation failed. Please check your HTML content.")
        # Tab for Unlock PDF
        elif option == "Unlock PDF":
            st.title("Unlock Protected PDF")

            # File uploader for the protected PDF
            uploaded_pdf = st.file_uploader("Upload a Protected PDF", type="pdf")

            if uploaded_pdf is not None:
                st.success("PDF uploaded successfully!")

                # Password input for unlocking the PDF
                password = st.text_input("Enter the PDF password", type="password")

                if password:
                    # Unlock the PDF
                    unlocked_pdf_stream = unlock_pdf(uploaded_pdf, password)

                    if unlocked_pdf_stream:
                        # Provide a download button for the unlocked PDF
                        st.download_button(
                            label="Download Unlocked PDF",
                            data=unlocked_pdf_stream,
                            file_name="unlocked.pdf",
                            mime="application/pdf"
                        )

        

        # Tab for Protect PDF
        elif option == "Protect PDF":
            st.subheader("Protect PDF")
            uploaded_file = st.file_uploader("Choose a PDF file to protect", type="pdf", key="protect_pdf")
            password = st.text_input("Enter password to protect PDF", type="password")

            if st.button("Protect PDF"):
                if uploaded_file and password:
                    protected_pdf = protect_pdf(uploaded_file, password)
                    if protected_pdf is not None:
                        st.success("PDF protected successfully!")
                        st.download_button("Download Protected PDF", protected_pdf, "protected.pdf")
                    else:
                        st.error("Failed to protect the PDF. Please try again.")
                else:
                    st.warning("Please upload a PDF file and enter a password.")

        # Tab for Organize PDF
        elif option == "Organize PDF":
            st.title("PDF Reorganizer")

            uploaded_pdf = st.file_uploader("Upload PDF", type=["pdf"])
            if uploaded_pdf is not None:
                st.success("PDF uploaded successfully!")
                
                if st.button("Reorganize PDF"):
                    organized_content = reorganize_pdf(uploaded_pdf)
                    
                    if organized_content:
                        st.download_button(
                            "Download Reorganized PDF",
                            organized_content,
                            "reorganized.pdf",
                            mime="application/pdf"
                        )
                    else:
                        st.error("Failed to reorganize the PDF. Please check the input and try again.")
        # Tab for Repair PDF
        elif option == "Repair PDF":
                st.title("PDF Repair Tool")

                uploaded_pdf = st.file_uploader("Upload PDF", type=["pdf"])
                if uploaded_pdf is not None:
                    st.success("PDF uploaded successfully!")
                    
                    if st.button("Repair PDF"):
                        repaired_content = repair_pdf(uploaded_pdf)
                        
                        if repaired_content:
                            st.download_button(
                                "Download Repaired PDF",
                                repaired_content,
                                "repaired.pdf",
                                mime="application/pdf"
                            )
                        else:
                            st.error("Failed to repair the PDF. Please check the input and try again.")


    


        # File Uploads and Conversion Process
        elif option == "PDF to JPG":
            st.title("PDF to JPG Converter")
            pdf_file = st.file_uploader("Upload PDF File", type=["pdf"])
            if pdf_file:
                images = convert_pdf_to_jpg(pdf_file)
                if images:
                    st.write("Converted Images:")
                    for i, img in enumerate(images):
                        st.image(img, caption=f"Page {i+1}", use_column_width=True)
                        img_bytes = io.BytesIO()
                        img.save(img_bytes, format="JPEG")
                        st.download_button(label=f"Download Page {i+1} as JPG",
                                        data=img_bytes.getvalue(),
                                        file_name=f"page_{i+1}.jpg",
                                        mime="image/jpeg")
                else:
                    st.warning("No images found in the PDF.")
                st.success("Conversion Successful!")

        elif option == "JPG to PDF":
            st.title("JPG to PDF Converter")
            image_files = st.file_uploader("Upload JPG Files", type=["jpg", "jpeg", "png"], accept_multiple_files=True)
            if image_files:
                pdf_bytes = convert_jpg_to_pdf(image_files)
                st.download_button("Download PDF", data=pdf_bytes, file_name="converted.pdf", mime="application/pdf")
                st.success("Conversion Successful!")
            else:
                st.warning("Please upload at least one JPG file.")

        elif option == "Excel to PDF":
            excel_file = st.file_uploader("Upload Excel File", type=["xlsx", "xls"])
            if excel_file:
                pdf_buffer = convert_excel_to_pdf(excel_file)
                st.download_button("Download PDF", data=pdf_buffer, file_name="converted.pdf", mime="application/pdf")
                st.success("Conversion Successful!")

        elif option == "PDF to Excel":
            st.subheader("PDF to Excel")
            uploaded_file = st.file_uploader("Upload PDF to convert to Excel", type=["pdf"])
            
            if st.button("Convert to Excel"):
                if uploaded_file:
                    excel_file = convert_pdf_to_excel(uploaded_file)
                    st.success("PDF converted to Excel successfully!")
                    
                    # Add the download button
                    st.download_button(
                        label="Download Excel File",
                        data=excel_file,
                        file_name="converted_excel.xlsx",
                        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                    )
                else:
                    st.warning("Please upload a PDF file to convert.")

        elif option == "PowerPoint to PDF":
            st.subheader("PowerPoint to PDF")
            uploaded_file = st.file_uploader("Upload a PowerPoint file (.pptx)", type="pptx")
            if uploaded_file is not None:
                pdf_file = convert_powerpoint_to_pdf(uploaded_file)
                if pdf_file is not None:
                    st.download_button(
                        label="Download PDF",
                        data=pdf_file,  # Ensure pdf_file is valid binary data
                        file_name="converted_file.pdf",
                        mime="application/pdf"
                    )
                else:
                    st.error("Failed to generate PDF. Please check the input file.")


        elif option == "PDF to PowerPoint":
            st.subheader("PDF to PowerPoint")
            uploaded_file = st.file_uploader("Upload PDF to convert to PowerPoint", type=["pdf"])
            
            if st.button("Convert to PowerPoint"):
                if uploaded_file:
                    powerpoint_file = convert_pdf_to_powerpoint(uploaded_file)
                    st.success("PDF converted to PowerPoint successfully!")
                    
                    # Add the download button
                    st.download_button(
                        label="Download PowerPoint File",
                        data=powerpoint_file,
                        file_name="converted_presentation.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    )
                else:
                    st.warning("Please upload a PDF file to convert.")

        elif option == "Word to PDF":
            
            uploaded_file = st.file_uploader("Upload a DOCX file", type="docx")

            if uploaded_file is not None:
                try:
                    document = Document(uploaded_file)
                    st.success("DOCX file loaded successfully.")

                    # Convert DOCX to PDF
                    pdf_buffer = convert_word_to_pdf(uploaded_file)

                    # Provide download button for the converted PDF
                    st.download_button(
                        label="Download PDF",
                        data=pdf_buffer,
                        file_name="converted.pdf",
                        mime="application/pdf"
                    )

                except BadZipFile:
                    st.error("The uploaded file is not a valid DOCX file. Please check the file and try again.")
                except UnicodeEncodeError as e:
                    st.error(f"Encoding error: {e}. Please ensure the document contains valid characters.")
                except Exception as e:
                    st.error(f"An unexpected error occurred: {e}")
        elif option == "PDF to Word":
            uploaded_pdf = st.file_uploader("Upload a PDF file", type="pdf")

            if uploaded_pdf is not None:
                word_file = convert_pdf_to_word(uploaded_pdf)
                
                if word_file is not None:
                    st.success("Conversion successful!")
                    
                    # Providing the download button
                    st.download_button(
                        label="Download Word Document",
                        data=word_file,
                        file_name="converted_file.docx",
                        mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                    )
                else:
                    st.error("Failed to convert PDF to Word.")


    if choice == "Image Features":
        st.title("Image Features")
        st.title("Welcome to the Image Editor Tools :")

      

        image_tabs = st.tabs(["Convert Image", "Compress Image", "Remove Background", "Resize Image & Passport Photo Options", "Image Watermarking Tool"])

        # Convert Image Tab
        with image_tabs[0]:
            st.header("Convert Image")
        
            uploaded_image = st.file_uploader("Upload an Image", type=["jpg", "jpeg", "png", "gif", "bmp", "tiff"], key="converter_uploader")
            
            if uploaded_image:
                image = Image.open(uploaded_image)
                
                # Conversion options
                output_format = st.selectbox("Select Output Format", ["JPEG", "PNG", "GIF", "BMP", "TIFF"], key="format_select")
                quality = st.slider("Quality (1-100)", min_value=1, max_value=100, value=85, step=1, key="quality_slider")
                
                if st.button("Convert Image"):
                    converted_image = convert_image(image, output_format, quality if output_format == "JPEG" else None)
                    st.image(converted_image, caption=f"Converted Image ({output_format})", use_column_width=True)
                    st.download_button("Download Image", converted_image, file_name=f"converted_image.{output_format.lower()}", mime=f"image/{output_format.lower()}")

        # Compress Image Tab
        with image_tabs[1]:
            st.header("Compress Image")
            uploaded_image_compress = st.file_uploader("Upload an Image for Compression", type=["jpg", "jpeg", "png"], key="compress_image_uploader")

            if uploaded_image_compress:
                compression_level = st.selectbox("Select Compression Level", ["Low", "Medium", "High"], key="compress_image_level")
                quality = st.selectbox("Select Quality", ["Low", "Medium", "High"], key="compress_image_quality")
                resolution = st.selectbox("Select Resolution", ["Original", "Full HD", "HD", "SD"], key="compress_image_resolution")
                
                if st.button("Compress"):
                    compressed_image, quality, compressed_image_size_kb = compress_image(uploaded_image_compress, compression_level, quality, resolution)
                    
                    st.image(compressed_image, caption=f"Compressed Image (Quality: {quality})", use_column_width=True)
                    
                    st.write(f"Estimated Compressed Image Size: {compressed_image_size_kb:.2f} KB")
                    
                    st.download_button("Download Compressed Image", compressed_image, file_name="compressed_image.jpg", mime="image/jpeg")

        # Remove Background Tab
        with image_tabs[2]:
            st.header("Remove Background")
            uploaded_image = st.file_uploader("Upload an Image", type=["jpg", "jpeg", "png"])
            if uploaded_image is not None:
                st.image(uploaded_image, caption="Uploaded Image", use_column_width=True)
                background_color = st.color_picker("Pick a Background Color", "#FFFFFF")  # Color picker for background color
                transparency = st.slider("Select Transparency Level", 0.0, 1.0, 1.0)  # Slider for transparency

                # Display the updated image dynamically as the user adjusts transparency
                if uploaded_image:
                    removed_bg_image = remove_background(uploaded_image, background_color, transparency)
                    st.image(removed_bg_image, caption="Background Removed", use_column_width=True)

                # Allow the user to download the final image
                buffer = io.BytesIO()
                removed_bg_image.save(buffer, format="PNG")  # Save as PNG to preserve transparency
                st.download_button("Download Image", data=buffer.getvalue(), file_name="background_removed.png", mime="image/png")

       

       
        # Resize Image & Passport Photo Options Tab
        with image_tabs[3]:
            st.header("Resize Image & Passport Photo Options")
            uploaded_image = st.file_uploader("Upload an image to resize", type=["jpg", "jpeg", "png"], key="resize_image")
            
            if uploaded_image is not None:
                resize_type = st.radio("Select Resize Option", ["Custom Resize", "Multiple Passport Photos"], key="resize_type")
                resolution = st.selectbox("Select resolution", ["Original", "Full HD", "HD", "SD"], key="resize_resolution")

                if resize_type == "Custom Resize":
                    width = st.number_input("Width (px)", min_value=1, step=1, key="custom_width")
                    height = st.number_input("Height (px)", min_value=1, step=1, key="custom_height")
                    if st.button("Resize Image", key="resize_button"):
                        img = Image.open(uploaded_image).resize((width, height))
                        buffer = BytesIO()
                        img.save(buffer, format="JPEG")
                        buffer.seek(0)
                        st.download_button("Download Resized Image", data=buffer, file_name="resized_image.jpg", mime="image/jpeg")
                
                elif resize_type == "Multiple Passport Photos":
                    photo_size_option = st.selectbox("Select Photo Size", ["College ID (35mm x 45mm)", "Passport (51mm x 51mm)"], key="photo_size")
                    number_of_photos = st.slider("Number of Photos", min_value=1, max_value=20, value=20)
                    add_border = st.checkbox("Add Border", value=False)
                    border_color = st.color_picker("Border Color", "#000000")
                    border_thickness = st.slider("Border Thickness", 1, 10, 2)
                    add_dob = st.checkbox("Add Date of Birth (DOB)", value=False)
                    dob = st.text_input("Enter Date of Birth (if applicable)")

                    if st.button("Generate Photos"):
                        # Define photo sizes in pixels (assuming 300 DPI)
                        dpi = 300
                        photo_size_mapping = {
                            "College ID (35mm x 45mm)": (int(35 * dpi / 25.4), int(45 * dpi / 25.4)),
                            "Passport (51mm x 51mm)": (int(51 * dpi / 25.4), int(51 * dpi / 25.4))
                        }
                        
                        if photo_size_option in photo_size_mapping:
                            photo_width, photo_height = photo_size_mapping[photo_size_option]

                            # A4 dimensions in pixels
                            a4_width, a4_height = 2480, 3508

                            # Calculate the number of rows and columns
                            cols = a4_width // (photo_width + 10)  # 10 pixels margin
                            rows = a4_height // (photo_height + 10)

                            # Create a blank A4 page
                            a4_page = Image.new("RGB", (a4_width, a4_height), "white")
                            
                            # Load and resize the uploaded image
                            photo_image = Image.open(uploaded_image).resize((photo_width, photo_height))

                            if add_border:
                                # Create a new image with border
                                bordered_image = Image.new("RGB", (photo_width + 2 * border_thickness, photo_height + 2 * border_thickness), border_color)
                                bordered_image.paste(photo_image, (border_thickness, border_thickness))
                                photo_image = bordered_image
                            
                            # Add DOB if needed
                            if add_dob and dob:
                                draw = ImageDraw.Draw(photo_image)
                                try:
                                    # Use a default font or a user-provided one
                                    font = ImageFont.truetype("arial.ttf", 24)
                                except IOError:
                                    font = ImageFont.load_default()

                                # Calculate text width and height
                                text_width, text_height = draw.textbbox((0, 0), dob, font=font)[2:4]
                                
                                text_x = (photo_width - text_width) / 2
                                text_y = photo_height - text_height - 10  # Position the text at the bottom
                                draw.text((text_x, text_y), dob, fill="black", font=font)
                            
                            # Paste the photos onto the A4 page
                            for i in range(number_of_photos):
                                x = (i % cols) * (photo_width + 10) + 10
                                y = (i // cols) * (photo_height + 10) + 10
                                if y + photo_height > a4_height:
                                    break
                                a4_page.paste(photo_image, (x, y))
                            
                            buffer = BytesIO()
                            a4_page.save(buffer, format="JPEG")
                            buffer.seek(0)
                            st.download_button("Download A4 Page with Multiple Photos", data=buffer, file_name="multiple_photos.jpg", mime="image/jpeg")
                        else:
                            st.error("Selected photo size option is not available.")

        
        # Image Watermarking Tool Tab
        with image_tabs[4]:
            st.title("Image Watermarking Tool")

            # Upload images
            uploaded_files = st.file_uploader("Upload Images", accept_multiple_files=True, type=["png", "jpg", "jpeg"])

            # Select watermark type
            watermark_type = st.selectbox("Select Watermark Type", ["Text", "Image"])

            # If text watermark is selected
            if watermark_type == "Text":
                text = st.text_input("Enter Watermark Text")
                font_size = st.slider("Font Size", 10, 200, 50)
                color = st.color_picker("Text Color", "#FFFFFF")
                opacity = st.slider("Opacity", 0, 255, 128)
                rotation = st.slider("Rotation", 0, 360, 0)
                
                # Select text position
                position = st.selectbox("Select Watermark Position", ["Top-Left", "Top-Right", "Center", "Bottom-Left", "Bottom-Right"])

                # Load font
                try:
                    font = ImageFont.truetype("arial.ttf", font_size)
                except IOError:
                    font = ImageFont.load_default()

            # If image watermark is selected
            elif watermark_type == "Image":
                watermark_image_file = st.file_uploader("Upload Watermark Image", type=["png", "jpg", "jpeg"])
                opacity = st.slider("Opacity", 0, 255, 128)
                scale = st.slider("Scale", 0.1, 1.0, 0.5)
                
                # Select image watermark position
                position = st.selectbox("Select Watermark Position", ["Top-Left", "Top-Right", "Center", "Bottom-Left", "Bottom-Right"])

            # Process the images
            if st.button("Apply Watermark"):
                if uploaded_files:
                    for uploaded_file in uploaded_files:
                        # Open the uploaded image
                        image = Image.open(uploaded_file).convert("RGBA")
                        
                        # Apply text watermark
                        if watermark_type == "Text":
                            # Create an overlay image for the text
                            text_overlay = Image.new("RGBA", image.size, (0, 0, 0, 0))
                            draw = ImageDraw.Draw(text_overlay)
                            
                            # Get text bounding box
                            bbox = draw.textbbox((0, 0), text, font=font)
                            text_width = bbox[2] - bbox[0]
                            text_height = bbox[3] - bbox[1]

                            text_x, text_y = {
                                "Top-Left": (10, 10),
                                "Top-Right": (image.width - text_width - 10, 10),
                                "Center": ((image.width - text_width) // 2, (image.height - text_height) // 2),
                                "Bottom-Left": (10, image.height - text_height - 10),
                                "Bottom-Right": (image.width - text_width - 10, image.height - text_height - 10)
                            }[position]

                            # Draw text on the overlay
                            draw.text((text_x, text_y), text, font=font, fill=(ImageColor.getrgb(color) + (opacity,)))
                            
                            # Composite the text overlay with the original image
                            watermarked_image = Image.alpha_composite(image, text_overlay)

                        # Apply image watermark
                        elif watermark_type == "Image" and watermark_image_file:
                            watermark_image = Image.open(watermark_image_file).convert("RGBA")
                            watermark_image = watermark_image.resize((int(watermark_image.width * scale), int(watermark_image.height * scale)))
                            watermark_image.putalpha(opacity)
                            
                            # Create an overlay for watermark
                            watermark_overlay = Image.new("RGBA", image.size)
                            position = {
                                "Top-Left": (10, 10),
                                "Top-Right": (image.width - watermark_image.width - 10, 10),
                                "Center": ((image.width - watermark_image.width) // 2, (image.height - watermark_image.height) // 2),
                                "Bottom-Left": (10, image.height - watermark_image.height - 10),
                                "Bottom-Right": (image.width - watermark_image.width - 10, image.height - watermark_image.height - 10)
                            }[position]
                            watermark_overlay.paste(watermark_image, position, watermark_image)
                            watermarked_image = Image.alpha_composite(image, watermark_overlay)

                        # Save the watermarked image
                        output_path = io.BytesIO()
                        watermarked_image.save(output_path, format="PNG")
                        output_path.seek(0)

                        # Display the watermarked image
                        st.image(watermarked_image, caption=f"Watermarked {uploaded_file.name}")

                        # Provide download option
                        st.download_button("Download Watermarked Image", data=output_path, file_name=f"watermarked_{uploaded_file.name}", mime="image/png")
                else:
                    st.error("No images uploaded.")
           



    elif choice == "Video Features":
        video_file = None
        
        st.title("Video Features")
        st.title("Welcome to the Video Compressor and Trimmer Tool :")

        st.header("Editing Tools")
        st.write("Select an editing tool below:")
        feature = st.selectbox("Choose a feature", ["Compress Video", "Trim Video"])

        

        if feature == "Compress Video":
           st.title("Compress Video")
           video_file = st.file_uploader("Upload a video", type=["mp4", "avi", "mov"])

        if video_file:
                st.video(video_file)
                compression_level = st.slider("Select compression level (bitrate in kbps)", 100, 5000, 1500)
                output_format = st.selectbox("Select output format", ["MP4", "AVI", "MOV"])
                
                if st.button("Compress"):
                    compressed_video = compress_video(video_file, compression_level, output_format)
                    st.success("Video compressed successfully!")
                    st.download_button(
                        label="Download Compressed Video",
                        data=compressed_video,
                        file_name=f"compressed_video.{output_format.lower()}",
                        mime="video/mp4"
                    )

        elif feature == "Trim Video":
            st.header("Trim Video Feature")
            video_file = st.file_uploader("Upload a video for trimming", type=["mp4", "mov", "avi"])

            if video_file is not None:
                st.video(video_file)
                start_time = st.number_input("Enter start time (in seconds)", min_value=0)
                end_time = st.number_input("Enter end time (in seconds)", min_value=0)

                if st.button("Trim Video"):
                    trimmed_video_path = trim_video(video_file, start_time, end_time)
                    if trimmed_video_path:
                        st.video(trimmed_video_path)
                        st.success("Video trimmed successfully!")
                        # Provide download button for the trimmed video
                        with open(trimmed_video_path, "rb") as f:
                            video_bytes = f.read()
                        st.download_button(
                            label="Download Trimmed Video",
                            data=video_bytes,
                            file_name=f"trimmed_{video_file.name}",
                            mime="video/mp4"
                        )
                    else:
                        st.error("Error trimming the video.")
            else:
                st.warning("Please upload a video file.")
                        
if __name__ == "__main__":
    main()
